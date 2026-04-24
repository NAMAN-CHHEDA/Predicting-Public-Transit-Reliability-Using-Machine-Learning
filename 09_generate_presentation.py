import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REPORT_DIR = ROOT / "project_outputs"
FIG_DIR = REPORT_DIR / "figures"
TABLE_DIR = REPORT_DIR / "tables"
SLIDE_PATH = ROOT / "Transit_Reliability_Project_Presentation.pptx"


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Predicting Public Transit Reliability"
    slide.placeholders[1].text = "Data 245 In-class Project Presentation\nTeam: [Add team names]"


def add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.4), width=Inches(8.2))
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(8.4), Inches(0.5))
    textbox.text_frame.text = caption


def add_metrics_slide(prs: Presentation, metrics_df: pd.DataFrame) -> None:
    best = metrics_df.sort_values("f1", ascending=False).iloc[0]
    add_bullets(
        prs,
        "Model Results (Multiple Metrics)",
        [
            f"Best model: {best['model']}",
            f"F1: {best['f1']:.3f}, ROC-AUC: {best['roc_auc']:.3f}, PR-AUC: {best['pr_auc']:.3f}",
            "Compared Logistic Regression, Random Forest, and Gradient Boosting",
            "Detailed metrics are in project_outputs/tables/model_metrics.csv",
        ],
    )


def add_appendix_rubric_slide(prs: Presentation) -> None:
    bullets = [
        "Code Walkthrough (5): 08_modeling_and_reporting.py and end-to-end pipeline scripts",
        "Presentation Skills/Time Mgmt (5): Use speaker notes and timed dry run",
        "Discussion/Q&A (5): Prepare trade-off answers for model selection and metrics",
        "Demo (5): Run model script + show saved model inference",
        "Visualization (2): Heatmap, target distribution, ROC/PR, confusion matrix",
        "Version Control (3): Add public GitHub repo URL in this slide before submission",
        "Lessons Learned (5): Included in dedicated lessons-learned slide",
        "Teamwork + Pair Programming (2): Add screenshots/logs in appendix",
        "Agile/Scrum (3): Add sprint board URL, meeting notes, backlog evidence",
        "Slides (5): This .pptx deck exported and submitted",
        "Saved model for quick demo (3): project_outputs/models/best_transit_delay_model.pkl",
        "Creative presentation techniques (2): Add transitions/animations before final export",
        "CRediT statement: Included in dedicated contributor roles slide",
        "Generative AI usage transparency: List tools and exactly where they were used",
    ]
    add_bullets(prs, "Appendix: Rubric Criteria Mapping", bullets)


def main() -> None:
    metrics_df = pd.read_csv(TABLE_DIR / "model_metrics.csv")
    with open(TABLE_DIR / "run_summary.json", "r", encoding="utf-8") as f:
        summary = json.load(f)

    prs = Presentation()

    add_title_slide(prs)
    add_bullets(
        prs,
        "Problem and Goal",
        [
            "Goal: Predict whether a transit stop event will be delayed",
            "Business value: Better rider communication and dispatch planning",
            f"Dataset size: {summary['dataset_rows']} rows and {summary['dataset_columns']} columns",
        ],
    )
    add_bullets(
        prs,
        "Data Pipeline and Feature Engineering",
        [
            "Merged schedule + realtime + weather data",
            "Engineered time, trip progression, weather, and speed features",
            "Balanced classes with SMOTE in feature engineering pipeline",
        ],
    )
    add_image_slide(
        prs,
        "Exploratory Data Analysis: Correlation Heatmap",
        FIG_DIR / "correlation_heatmap.png",
        "Heatmap used to identify feature relationships and potential multicollinearity.",
    )
    add_image_slide(
        prs,
        "Target Distribution",
        FIG_DIR / "target_distribution.png",
        "Balanced delay classes used for robust model training.",
    )
    add_metrics_slide(prs, metrics_df)
    add_image_slide(
        prs,
        "Model Quality: ROC Curves",
        FIG_DIR / "roc_curves.png",
        "ROC curves for all candidate models.",
    )
    add_image_slide(
        prs,
        "Model Quality: Precision-Recall Curves",
        FIG_DIR / "pr_curves.png",
        "PR curves emphasize positive-class performance.",
    )
    add_image_slide(
        prs,
        "Best Model Confusion Matrix",
        FIG_DIR / "confusion_matrix_best_model.png",
        "Classification behavior on held-out test data.",
    )

    feature_importance = FIG_DIR / "feature_importance_top10.png"
    if feature_importance.exists():
        add_image_slide(
            prs,
            "Feature Importance (Best Tree-Based Model)",
            feature_importance,
            "Top contributors to transit delay prediction.",
        )

    add_bullets(
        prs,
        "Code Walkthrough and Demo Plan",
        [
            "Step 1: Run 08_modeling_and_reporting.py to train/evaluate models",
            "Step 2: Show saved model file and quick inference script execution",
            "Step 3: Walk through key functions: preprocessing, metrics, and model selection",
        ],
    )
    add_bullets(
        prs,
        "Lessons Learned and Technical Difficulties",
        [
            "Realtime transit data quality required careful filtering and joins",
            "Class imbalance initially hurt recall before resampling",
            "Best practice: track multiple metrics, not only accuracy",
            "Future work: route-specific models and temporal cross-validation",
        ],
    )
    add_bullets(
        prs,
        "Teamwork, Agile, and Pair Programming Evidence",
        [
            "Add sprint board link (Trello/Jira) and sprint artifacts",
            "Add pair-programming screenshots or VS Code LiveShare evidence",
            "Add commit history snapshots from GitHub and role ownership notes",
        ],
    )
    add_bullets(
        prs,
        "CRediT Contributor Roles Statement",
        [
            "[Name 1] Conceptualization, Data Curation, and Investigation",
            "[Name 2] Methodology, Software, and Validation",
            "[Name 3] Visualization, Writing, and Presentation",
            "Replace placeholders with exact team member contributions",
        ],
    )
    add_appendix_rubric_slide(prs)
    add_bullets(
        prs,
        "Appendix: External Evidence Links",
        [
            "GitHub repo URL: [insert public URL]",
            "Sprint board URL: [insert Trello/Jira URL]",
            "Elevator pitch video URL: [insert link]",
            "Cloud model URL (if model >2MB): [insert link + file size]",
            "Generative AI tools used and where: [insert transparent disclosure]",
        ],
    )

    prs.save(SLIDE_PATH)
    print(f"Saved slide deck: {SLIDE_PATH}")


if __name__ == "__main__":
    main()
